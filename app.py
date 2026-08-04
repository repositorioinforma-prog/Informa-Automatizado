import io
import os
import tempfile
from io import BytesIO

import numpy as np
import openpyxl
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import pyreadstat
import requests
import scipy.stats as stats
import seaborn as sns
import streamlit as st
from PIL import Image
from adjustText import adjust_text
from matplotlib import pyplot as plt
from matplotlib.colors import LinearSegmentedColormap, Normalize
from openpyxl.styles import Font, PatternFill
from openpyxl.utils import get_column_letter
from scipy.stats import chi2_contingency, pearsonr
from sklearn.linear_model import LinearRegression
from sklearn.metrics import mean_squared_error, r2_score
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import LabelEncoder
from statsmodels.stats.outliers_influence import variance_inflation_factor
from statsmodels.tools.tools import add_constant

# Mantido apenas por compatibilidade, embora não esteja em uso na CA manual
import prince

# Sweetviz protegido para não derrubar o app inteiro
try:
    import sweetviz as sv
    SWEETVIZ_AVAILABLE = True
    SWEETVIZ_IMPORT_ERROR = ""
except Exception as e:
    sv = None
    SWEETVIZ_AVAILABLE = False
    SWEETVIZ_IMPORT_ERROR = str(e)

# python-pptx protegido, usado na exportação dos mapas para PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
    PPTX_IMPORT_ERROR = ""
except Exception as e:
    Presentation = None
    Inches = None
    Pt = None
    PPTX_AVAILABLE = False
    PPTX_IMPORT_ERROR = str(e)


# =========================
# Configuração da página
# =========================
# Variáveis normalmente usadas nos cruzamentos de Mapas de Correspondência.
# Quando um arquivo é carregado, o app verifica quais delas existem e já
# sugere/pré-seleciona no módulo "Múltiplas Variáveis". Edite esta lista
# livremente conforme o questionário mudar.
VARIAVEIS_INTERESSE_PADRAO = [
    "P1", "P2", "P3", "P4", "P5", "juncaoreligiao", "P7", "P7_C",
]

# Nomes amigáveis para exibir no lugar do código da variável (ex.: nos
# títulos dos slides do PowerPoint exportado). Variáveis sem entrada aqui
# continuam aparecendo com o próprio código (ex.: "juncaoreligiao").
ROTULOS_VARIAVEIS = {
    "P1": "SEXO",
    "P2": "IDADE",
    "P3": "RENDA",
    "P4": "ESCOLARIDADE",
    "P5": "RAÇA",
    "P7": "REGIÕES",
    "P7_C": "CAPITAL",
}


def _rotulo_amigavel(variavel):
    """Retorna o nome amigável da variável, se houver; senão, o próprio código."""
    return ROTULOS_VARIAVEIS.get(variavel, variavel)


st.set_page_config(page_title="Análise de Dados Automática", layout="wide")

st.markdown(
    """
    <style>
    :root {
        --accent: #2563eb;
        --accent-hover: #1d4ed8;
        --text-main: #111827;
        --text-muted: #4b5563;
        --border: #94a3b8;
        --bg-card: #f8fafc;
        --bg-app: #ffffff;
        --sidebar-bg: #1f2937;
    }

    html, body, .stApp {
        background: var(--bg-app);
        color: var(--text-main);
        font-size: 16px;
    }

    /* texto geral */
    .stApp, .stApp p, .stApp label, .stApp span, .stApp div,
    .stMarkdown, .stText, h1, h2, h3, h4, h5, h6 {
        color: var(--text-main);
    }

    h1, h2, h3 {
        font-weight: 700;
        border-bottom: 3px solid var(--accent);
        padding-bottom: 6px;
        display: inline-block;
    }

    /* sidebar escura */
    [data-testid="stSidebar"] {
        background-color: var(--sidebar-bg);
        color: #f9fafb;
    }

    /* Só força cor clara em elementos de texto "folha" (p, label, span,
       small). NÃO inclui 'div' de propósito: divs também envolvem
       componentes como selectbox/multiselect, que têm fundo BRANCO — se
       forçássemos texto claro neles, o texto ficaria branco sobre
       branco (invisível). Esses componentes recebem sua própria regra
       de cor mais abaixo, com especificidade maior. */
    [data-testid="stSidebar"] p,
    [data-testid="stSidebar"] label,
    [data-testid="stSidebar"] span,
    [data-testid="stSidebar"] small {
        color: #f9fafb !important;
    }

    /* selects/multiselects dentro da sidebar têm fundo branco (definido
       mais abaixo), então o texto deles precisa continuar escuro —
       especificidade maior que a regra acima para garantir a prioridade */
    [data-testid="stSidebar"] [data-baseweb="select"],
    [data-testid="stSidebar"] [data-baseweb="select"] * {
        color: var(--text-main) !important;
    }

    /* file uploader */
    [data-testid="stFileUploaderDropzone"] {
        background-color: #111827 !important;
        border: 2px dashed #6b7280 !important;
        border-radius: 10px;
    }

    [data-testid="stFileUploaderDropzone"] * {
        color: #f9fafb !important;
    }

    [data-testid="stFileUploaderDropzone"] button {
        background-color: var(--accent) !important;
        color: #ffffff !important;
        border: none !important;
        border-radius: 8px;
    }

    /* inputs de texto e número */
    textarea, input {
        background-color: #ffffff !important;
        color: var(--text-main) !important;
        border: 1.5px solid var(--border) !important;
        border-radius: 6px;
    }

    textarea:focus, input:focus {
        border-color: var(--accent) !important;
        outline: none !important;
        box-shadow: 0 0 0 2px rgba(37, 99, 235, 0.25) !important;
    }

    /* selects e multiselects (componentes BaseWeb, não são <input> puro) */
    [data-baseweb="select"] > div {
        background-color: #ffffff !important;
        border: 1.5px solid var(--border) !important;
        border-radius: 6px !important;
    }

    [data-baseweb="select"] * {
        color: var(--text-main) !important;
    }

    [data-baseweb="tag"] {
        background-color: var(--accent) !important;
        color: #ffffff !important;
        border-radius: 4px !important;
    }

    [data-baseweb="tag"] svg {
        fill: #ffffff !important;
    }

    [data-baseweb="popover"] [data-baseweb="menu"] {
        background-color: #ffffff !important;
        border: 1px solid var(--border) !important;
    }

    [data-baseweb="menu"] li {
        color: var(--text-main) !important;
    }

    [data-baseweb="menu"] li:hover {
        background-color: #eff6ff !important;
    }

    /* checkboxes e radios dentro da sidebar (fundo escuro) */
    [data-testid="stSidebar"] [data-testid="stCheckbox"] label,
    [data-testid="stSidebar"] [data-testid="stRadio"] label {
        color: #f9fafb !important;
    }

    /* tabelas */
    .dataframe {
        background-color: #ffffff;
        border: 1px solid var(--border);
        color: var(--text-main);
        border-radius: 8px;
    }

    .dataframe tbody tr:nth-child(odd) {
        background-color: #f1f5f9;
    }

    .dataframe tbody tr:nth-child(even) {
        background-color: #ffffff;
    }

    .dataframe thead {
        background-color: var(--accent);
        color: white;
    }

    /* botões (comuns, de formulário e de download) */
    .stButton button,
    .stFormSubmitButton button,
    .stDownloadButton button {
        background-color: var(--accent);
        border: none;
        border-radius: 8px;
        color: #ffffff;
        padding: 8px 18px;
        font-weight: 600;
        cursor: pointer;
        transition: background-color 0.15s ease-in-out;
    }

    .stButton button:hover,
    .stFormSubmitButton button:hover,
    .stDownloadButton button:hover {
        background-color: var(--accent-hover);
        color: #ffffff;
    }

    /* expanders viram "cards", ajuda a distinguir vários mapas na tela */
    [data-testid="stExpander"] {
        border: 1px solid var(--border) !important;
        border-radius: 10px !important;
        background-color: var(--bg-card) !important;
        margin-bottom: 14px;
    }

    [data-testid="stExpander"] summary {
        font-weight: 600;
        font-size: 1.05rem;
    }

    hr {
        border-top: 1.5px solid var(--border);
    }
    </style>
    """,
    unsafe_allow_html=True
)

# =========================
# Logo
# =========================
url = "https://institutoinforma.com.br/wp-content/uploads/2025/01/logo_informa.webp"
headers = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/110.0.0.0 Safari/537.36"
    )
}

try:
    response = requests.get(url, headers=headers, timeout=20)
    response.raise_for_status()
    image_bytes = BytesIO(response.content)
    logo = Image.open(image_bytes)
    st.image(logo, width=300)
except Exception:
    st.warning("Não foi possível carregar a logo do sistema.")


# =========================
# Funções utilitárias
# =========================
def carregar_dados(uploaded_file):
    """
    Carrega dados de CSV, Excel ou SPSS (SAV), preservando valores numéricos
    e associando labels apenas para exibição.
    """
    try:
        if uploaded_file.name.endswith(".csv"):
            try:
                dados = pd.read_csv(uploaded_file, encoding="utf-8")
            except UnicodeDecodeError:
                st.warning("Codificação 'utf-8' falhou. Tentando com 'latin1'.")
                dados = pd.read_csv(uploaded_file, encoding="latin1")

            st.success(f"Dados carregados com sucesso! Número de registros: {dados.shape[0]}")
            st.write("Visualização dos dados:", dados.head())
            return dados, dados

        elif uploaded_file.name.endswith(".xlsx"):
            dados = pd.read_excel(uploaded_file)
            st.success(f"Dados carregados com sucesso! Número de registros: {dados.shape[0]}")
            st.write("Visualização dos dados:", dados.head())
            return dados, dados

        elif uploaded_file.name.endswith(".sav"):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".sav") as tmp:
                tmp.write(uploaded_file.getvalue())
                tmp_path = tmp.name

            dados, meta = pyreadstat.read_sav(tmp_path)
            os.unlink(tmp_path)

            label_dict = {}
            for var in meta.variable_value_labels:
                label_dict[var] = {
                    int(k): v for k, v in meta.variable_value_labels[var].items()
                }

            dados_exibicao = dados.copy()
            for col in label_dict.keys():
                if col in dados.columns:
                    dados_exibicao[col] = (
                        dados[col].map(label_dict[col]).fillna(dados[col])
                    )

            st.success(f"Dados carregados com sucesso! Número de registros: {dados.shape[0]}")
            st.write("Visualização com labels:", dados_exibicao.head())
            st.write("Dados numéricos usados nos cálculos:", dados.head())

            return dados, dados_exibicao

        else:
            st.error("Formato de arquivo não suportado. Use CSV, Excel ou SPSS (.sav).")
            return None, None

    except Exception as e:
        st.error(f"Erro ao carregar dados: {e}")
        return None, None


def validar_arquivo(arquivo):
    """
    Função para validar se o arquivo carregado contém dados válidos.
    """
    if arquivo is not None:
        try:
            dados = pd.read_csv(arquivo)
            if dados.empty:
                st.error("O arquivo está vazio. Por favor, carregue um arquivo válido.")
                return None
            return dados
        except Exception as e:
            st.error(f"Erro ao ler o arquivo: {e}")
            return None
    else:
        st.warning("Nenhum arquivo carregado.")
        return None


def preprocessar_dados(dados):
    for coluna in dados.columns:
        if dados[coluna].apply(
            lambda x: str(x).replace(".", "", 1).isdigit() if pd.notnull(x) else False
        ).all():
            dados[coluna] = pd.to_numeric(dados[coluna], errors="coerce")
        else:
            dados[coluna] = dados[coluna].astype(str)
    return dados

# =========================
# Análise de Correspondência
# =========================
def _ca_manual(tabela_contingencia, n_components=2):
    """
    Implementação manual de Análise de Correspondência via SVD.
    """
    N = tabela_contingencia.to_numpy(dtype=float)
    n = N.sum()
    if n <= 0:
        raise ValueError("Tabela de contingência vazia.")

    P = N / n
    r = P.sum(axis=1)
    c = P.sum(axis=0)

    row_mask = r > 0
    col_mask = c > 0

    if row_mask.sum() < 2 or col_mask.sum() < 2:
        raise ValueError("Linhas/colunas insuficientes com massa > 0 para a CA.")

    P_rc = P[np.ix_(row_mask, col_mask)]
    r_rc = r[row_mask]
    c_rc = c[col_mask]

    D_r_inv_sqrt = np.diag(1.0 / np.sqrt(r_rc))
    D_c_inv_sqrt = np.diag(1.0 / np.sqrt(c_rc))

    expected = np.outer(r_rc, c_rc)
    S = D_r_inv_sqrt @ (P_rc - expected) @ D_c_inv_sqrt

    U, singvals, VT = np.linalg.svd(S, full_matrices=False)
    eigenvalues = singvals ** 2

    F = D_r_inv_sqrt @ U @ np.diag(singvals)
    G = D_c_inv_sqrt @ VT.T @ np.diag(singvals)

    k = min(n_components, F.shape[1])
    F = F[:, :k]
    G = G[:, :k]
    eigenvalues = eigenvalues[:k]

    row_index = tabela_contingencia.index[row_mask]
    col_index = tabela_contingencia.columns[col_mask]
    cols = [f"Dim{i+1}" for i in range(k)]

    coordenadas_linhas = pd.DataFrame(F, index=row_index, columns=cols)
    coordenadas_colunas = pd.DataFrame(G, index=col_index, columns=cols)

    return coordenadas_linhas, coordenadas_colunas, eigenvalues


def analise_correspondencia(dados):
    st.header("Análise de Correspondência")

    colunas_categoricas = st.multiselect(
        "Selecione as colunas categóricas para a análise:",
        dados.columns
    )

    if len(colunas_categoricas) < 2:
        st.warning("Selecione pelo menos duas colunas categóricas para realizar a análise.")
        return

    tabela_contingencia = pd.crosstab(
        dados[colunas_categoricas[0]],
        dados[colunas_categoricas[1]],
        normalize=False
    )

    try:
        coordenadas_linhas, coordenadas_colunas, eigenvalues = _ca_manual(
            tabela_contingencia,
            n_components=2
        )
    except Exception as e:
        st.error(f"Não foi possível calcular a Análise de Correspondência: {e}")
        return

    coordenadas_linhas *= -1
    coordenadas_colunas *= -1

    if coordenadas_linhas.shape[1] < 2 or coordenadas_colunas.shape[1] < 2:
        st.error(
            "A análise de correspondência não conseguiu gerar duas dimensões. "
            "Verifique os dados selecionados."
        )
        return

    st.subheader("Inércia explicada (variância)")
    explained_inertia = eigenvalues / eigenvalues.sum()
    for i, valor in enumerate(explained_inertia):
        st.write(f"Dim {i+1}: {valor * 100:.2f}%")

    if "rotulos_linhas" not in st.session_state:
        st.session_state["rotulos_linhas"] = {
            i: str(i) for i in coordenadas_linhas.index
        }
    else:
        for i in coordenadas_linhas.index:
            st.session_state["rotulos_linhas"].setdefault(i, str(i))

    if "rotulos_colunas" not in st.session_state:
        st.session_state["rotulos_colunas"] = {
            i: str(i) for i in coordenadas_colunas.index
        }
    else:
        for i in coordenadas_colunas.index:
            st.session_state["rotulos_colunas"].setdefault(i, str(i))

    if "deslocamentos_linhas" not in st.session_state:
        st.session_state["deslocamentos_linhas"] = {
            i: (0.0, 0.0) for i in coordenadas_linhas.index
        }
    else:
        for i in coordenadas_linhas.index:
            st.session_state["deslocamentos_linhas"].setdefault(i, (0.0, 0.0))

    if "deslocamentos_colunas" not in st.session_state:
        st.session_state["deslocamentos_colunas"] = {
            i: (0.0, 0.0) for i in coordenadas_colunas.index
        }
    else:
        for i in coordenadas_colunas.index:
            st.session_state["deslocamentos_colunas"].setdefault(i, (0.0, 0.0))

    legenda_linhas = st.text_input("Legenda para Linhas:", "Linhas")
    legenda_colunas = st.text_input("Legenda para Colunas:", "Colunas")
    mostrar_legenda = st.checkbox("Mostrar legenda (Linhas/Colunas)", value=False)
    editar_rotulos = st.checkbox("Editar rótulos e deslocamentos")

    if editar_rotulos:
        st.subheader("Editar Rótulos e Posições")

        st.markdown("**Linhas**")
        for i in coordenadas_linhas.index:
            st.session_state["rotulos_linhas"][i] = st.text_input(
                f"Novo rótulo para linha '{i}':",
                value=st.session_state["rotulos_linhas"][i],
                key=f"rotulo_linha_{i}"
            )

            dx, dy = st.session_state["deslocamentos_linhas"][i]
            dx = st.number_input(
                f"Deslocamento X para linha '{i}':",
                min_value=-1.0,
                max_value=1.0,
                value=float(dx),
                step=0.01,
                format="%.2f",
                key=f"num_desloc_x_linha_{i}"
            )
            dy = st.number_input(
                f"Deslocamento Y para linha '{i}':",
                min_value=-1.0,
                max_value=1.0,
                value=float(dy),
                step=0.01,
                format="%.2f",
                key=f"num_desloc_y_linha_{i}"
            )
            st.session_state["deslocamentos_linhas"][i] = (dx, dy)

            if st.checkbox(f"Remover ponto e rótulo da linha '{i}'", key=f"remover_linha_{i}"):
                st.session_state["rotulos_linhas"][i] = None

        st.markdown("**Colunas**")
        for i in coordenadas_colunas.index:
            st.session_state["rotulos_colunas"][i] = st.text_input(
                f"Novo rótulo para coluna '{i}':",
                value=st.session_state["rotulos_colunas"][i],
                key=f"rotulo_coluna_{i}"
            )

            dx, dy = st.session_state["deslocamentos_colunas"][i]
            dx = st.number_input(
                f"Deslocamento X para coluna '{i}':",
                min_value=-1.0,
                max_value=1.0,
                value=float(dx),
                step=0.01,
                format="%.2f",
                key=f"num_desloc_x_coluna_{i}"
            )
            dy = st.number_input(
                f"Deslocamento Y para coluna '{i}':",
                min_value=-1.0,
                max_value=1.0,
                value=float(dy),
                step=0.01,
                format="%.2f",
                key=f"num_desloc_y_coluna_{i}"
            )
            st.session_state["deslocamentos_colunas"][i] = (dx, dy)

            if st.checkbox(f"Remover ponto e rótulo da coluna '{i}'", key=f"remover_coluna_{i}"):
                st.session_state["rotulos_colunas"][i] = None

    max_deslocamento = 0.05
    pontos_linhas = []
    pontos_colunas = []

    for i, row in coordenadas_linhas.iterrows():
        if st.session_state["rotulos_linhas"][i] is None:
            continue
        x = row.iloc[0]
        y = row.iloc[1]
        dx, dy = st.session_state["deslocamentos_linhas"][i]
        dx = max(-max_deslocamento, min(max_deslocamento, dx * (x / abs(x) if x != 0 else 1)))
        dy = max(-max_deslocamento, min(max_deslocamento, dy * (y / abs(y) if y != 0 else 1)))
        pontos_linhas.append({
            "x": x,
            "y": y,
            "dx": dx,
            "dy": dy,
            "label": st.session_state["rotulos_linhas"][i]
        })

    for i, row in coordenadas_colunas.iterrows():
        if st.session_state["rotulos_colunas"][i] is None:
            continue
        x = row.iloc[0]
        y = row.iloc[1]
        dx, dy = st.session_state["deslocamentos_colunas"][i]
        dx = max(-max_deslocamento, min(max_deslocamento, dx * (x / abs(x) if x != 0 else 1)))
        dy = max(-max_deslocamento, min(max_deslocamento, dy * (y / abs(y) if y != 0 else 1)))
        pontos_colunas.append({
            "x": x,
            "y": y,
            "dx": dx,
            "dy": dy,
            "label": st.session_state["rotulos_colunas"][i]
        })

    st.subheader("Mapa estático")

    fig, ax = plt.subplots(figsize=(8, 6))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    texts = []
    linha_legenda_plotada = False
    coluna_legenda_plotada = False

    for p in pontos_linhas:
        label_legenda = legenda_linhas if (mostrar_legenda and not linha_legenda_plotada) else ""
        ax.scatter(p["x"], p["y"], color="blue", marker="o", s=20, label=label_legenda)
        if mostrar_legenda and not linha_legenda_plotada:
            linha_legenda_plotada = True

        txt = ax.text(
            p["x"] + p["dx"],
            p["y"] + p["dy"],
            p["label"],
            color="blue",
            fontsize=8,
            ha="center",
            va="bottom",
            fontweight="bold",
        )
        texts.append(txt)

    for p in pontos_colunas:
        label_legenda = legenda_colunas if (mostrar_legenda and not coluna_legenda_plotada) else ""
        ax.scatter(p["x"], p["y"], color="red", marker="^", s=30, label=label_legenda)
        if mostrar_legenda and not coluna_legenda_plotada:
            coluna_legenda_plotada = True

        txt = ax.text(
            p["x"] + p["dx"],
            p["y"] + p["dy"],
            p["label"],
            color="red",
            fontsize=8,
            ha="center",
            va="bottom",
            fontweight="bold",
        )
        texts.append(txt)

    adjust_text(
        texts,
        arrowprops=None,
        force_text=(0.5, 1),
        force_points=(0.5, 1),
        expand_text=(1.2, 1.5),
        expand_points=(1.2, 1.5),
        lim=100,
    )

    ax.axhline(0, color="black", linestyle="--", linewidth=0.8)
    ax.axvline(0, color="black", linestyle="--", linewidth=0.8)
    ax.grid(color="lightgray", linestyle="--", linewidth=0.5, alpha=0.7)
    ax.set_aspect("auto")

    if mostrar_legenda and (linha_legenda_plotada or coluna_legenda_plotada):
        ax.legend(loc="upper right", frameon=False, fontsize=12)

    ax.set_title("Análise de Correspondência", fontsize=14)
    st.pyplot(fig)

    st.subheader("Mapa interativo")

    fig_int = go.Figure()

    fig_int.add_trace(go.Scatter(
        x=[p["x"] for p in pontos_linhas],
        y=[p["y"] for p in pontos_linhas],
        mode="markers",
        name=legenda_linhas,
        marker=dict(color="blue", symbol="circle", size=8),
        showlegend=mostrar_legenda and len(pontos_linhas) > 0
    ))

    fig_int.add_trace(go.Scatter(
        x=[p["x"] for p in pontos_colunas],
        y=[p["y"] for p in pontos_colunas],
        mode="markers",
        name=legenda_colunas,
        marker=dict(color="red", symbol="triangle-up", size=10),
        showlegend=mostrar_legenda and len(pontos_colunas) > 0
    ))

    annotations = []
    for p in pontos_linhas:
        annotations.append(dict(
            x=p["x"] + p["dx"],
            y=p["y"] + p["dy"],
            xref="x",
            yref="y",
            text=p["label"],
            showarrow=False,
            font=dict(color="blue", size=10),
            xanchor="center",
            yanchor="bottom",
        ))
    for p in pontos_colunas:
        annotations.append(dict(
            x=p["x"] + p["dx"],
            y=p["y"] + p["dy"],
            xref="x",
            yref="y",
            text=p["label"],
            showarrow=False,
            font=dict(color="red", size=10),
            xanchor="center",
            yanchor="bottom",
        ))

    xs = [p["x"] for p in pontos_linhas] + [p["x"] for p in pontos_colunas]
    ys = [p["y"] for p in pontos_linhas] + [p["y"] for p in pontos_colunas]

    if xs and ys:
        x_min, x_max = min(xs), max(xs)
        y_min, y_max = min(ys), max(ys)
        dx = (x_max - x_min) * 0.1 or 0.1
        dy = (y_max - y_min) * 0.1 or 0.1
        x_range = [x_min - dx, x_max + dx]
        y_range = [y_min - dy, y_max + dy]
    else:
        x_range = [-1, 1]
        y_range = [-1, 1]

    fig_int.update_layout(
        template=None,
        title=dict(
            text="Análise de Correspondência",
            font=dict(color="black", size=14),
            x=0.5
        ),
        font=dict(color="black"),
        paper_bgcolor="white",
        plot_bgcolor="white",
        annotations=annotations,
        showlegend=mostrar_legenda,
        legend=dict(
            x=0.99, y=0.99,
            xanchor="right", yanchor="top",
            bgcolor="rgba(0,0,0,0)",
            font=dict(color="black", size=12),
        ),
        margin=dict(l=40, r=20, t=60, b=40),
        xaxis=dict(range=x_range),
        yaxis=dict(range=y_range),
    )

    fig_int.update_xaxes(
        title=dict(text="", font=dict(color="black")),
        showgrid=True,
        gridcolor="lightgray",
        gridwidth=0.5,
        griddash="dash",
        zeroline=False,
        tickfont=dict(color="black"),
        linecolor="black",
        mirror=True,
    )
    fig_int.update_yaxes(
        title=dict(text="", font=dict(color="black")),
        showgrid=True,
        gridcolor="lightgray",
        gridwidth=0.5,
        griddash="dash",
        zeroline=False,
        tickfont=dict(color="black"),
        linecolor="black",
        mirror=True,
    )

    fig_int.add_shape(
        type="line",
        x0=x_range[0], x1=x_range[1],
        y0=0, y1=0,
        line=dict(color="black", width=1, dash="dash"),
        xref="x", yref="y",
        layer="above"
    )
    fig_int.add_shape(
        type="line",
        x0=0, x1=0,
        y0=y_range[0], y1=y_range[1],
        line=dict(color="black", width=1, dash="dash"),
        xref="x", yref="y",
        layer="above"
    )

    config = {
        "editable": True,
        "edits": {"annotationPosition": True},
    }

    st.plotly_chart(fig_int, use_container_width=True, config=config)


# =========================================================
# Cálculo cacheado (crosstab + CA) — evita recalcular a cada rerun
# =========================================================
@st.cache_data(show_spinner=False)
def _crosstab_e_ca(dados, var_linha, var_coluna, n_components=2):
    """
    Monta a tabela de contingência e calcula a CA para o par de variáveis.
    Fica em cache: só recalcula se 'dados', 'var_linha' ou 'var_coluna'
    mudarem, então trocar um widget em um mapa não recalcula os outros.
    """
    tabela_contingencia = pd.crosstab(
        dados[var_linha],
        dados[var_coluna],
        normalize=False
    )
    coordenadas_linhas, coordenadas_colunas, eigenvalues = _ca_manual(
        tabela_contingencia,
        n_components=n_components
    )
    coordenadas_linhas = coordenadas_linhas * -1
    coordenadas_colunas = coordenadas_colunas * -1
    return coordenadas_linhas, coordenadas_colunas, eigenvalues


# =========================================================
# Estado compartilhado das categorias da variável principal
# =========================================================
def _preparar_estado_linhas_principal(dados, variavel_principal):
    """
    Prepara os dicionários de rótulos/deslocamentos das categorias da
    variável principal (linhas). Esse estado é compartilhado entre todos
    os cruzamentos gerados em `analise_correspondencia_multipla`, já que a
    variável principal é sempre a mesma nas linhas de cada mapa.
    """
    key_rotulos = f"rotulos_linhas_principal__{variavel_principal}"
    key_deslocamentos = f"deslocamentos_linhas_principal__{variavel_principal}"

    categorias = sorted(
        dados[variavel_principal].dropna().unique().tolist(), key=lambda v: str(v)
    )

    if key_rotulos not in st.session_state:
        st.session_state[key_rotulos] = {c: str(c) for c in categorias}
    else:
        for c in categorias:
            st.session_state[key_rotulos].setdefault(c, str(c))

    if key_deslocamentos not in st.session_state:
        st.session_state[key_deslocamentos] = {c: (0.0, 0.0) for c in categorias}
    else:
        for c in categorias:
            st.session_state[key_deslocamentos].setdefault(c, (0.0, 0.0))

    return key_rotulos, key_deslocamentos, categorias


def _configurar_categorias_principal(dados, variavel_principal):
    """
    Tela única para renomear, deslocar ou remover categorias da variável
    principal. Fica dentro de um formulário: só aplica (e só recalcula os
    mapas) quando o botão é clicado, não a cada tecla/clique.
    """
    key_rotulos, key_deslocamentos, categorias = _preparar_estado_linhas_principal(
        dados, variavel_principal
    )

    with st.expander(
        f"⚙️ Configurar categorias de '{variavel_principal}' "
        "— aplicado a todos os mapas abaixo",
        expanded=False
    ):
        st.caption(
            "Como a variável principal é sempre a mesma nas linhas de cada "
            "mapa, renomear, deslocar ou remover uma categoria aqui atualiza "
            "todos os cruzamentos de uma vez, ao clicar em Aplicar."
        )

        with st.form(key=f"form_principal__{variavel_principal}"):
            rotulo_tmp = {}
            dx_tmp = {}
            dy_tmp = {}
            remover_tmp = {}

            for c in categorias:
                col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
                rotulo_atual = st.session_state[key_rotulos][c]
                dx_atual, dy_atual = st.session_state[key_deslocamentos][c]

                with col1:
                    rotulo_tmp[c] = st.text_input(
                        f"Rótulo — '{c}'",
                        value=rotulo_atual if rotulo_atual is not None else str(c),
                        key=f"rotulo_principal__{variavel_principal}__{c}"
                    )
                with col2:
                    dx_tmp[c] = st.number_input(
                        "Δx", min_value=-1.0, max_value=1.0,
                        value=float(dx_atual), step=0.01, format="%.2f",
                        key=f"desloc_x_principal__{variavel_principal}__{c}"
                    )
                with col3:
                    dy_tmp[c] = st.number_input(
                        "Δy", min_value=-1.0, max_value=1.0,
                        value=float(dy_atual), step=0.01, format="%.2f",
                        key=f"desloc_y_principal__{variavel_principal}__{c}"
                    )
                with col4:
                    remover_tmp[c] = st.checkbox(
                        "Remover",
                        value=(rotulo_atual is None),
                        key=f"remover_principal__{variavel_principal}__{c}"
                    )

            aplicar = st.form_submit_button("Aplicar em todos os mapas")

            if aplicar:
                for c in categorias:
                    st.session_state[key_rotulos][c] = (
                        None if remover_tmp[c] else rotulo_tmp[c]
                    )
                    st.session_state[key_deslocamentos][c] = (dx_tmp[c], dy_tmp[c])
                st.success("Configuração aplicada a todos os cruzamentos.")

    return key_rotulos, key_deslocamentos


# =========================================================
# Mapa de Correspondência auxiliar (uso interno, parametrizado)
# =========================================================
def _renderizar_mapa_correspondencia(
    dados, var_linha, var_coluna, key_prefix, key_rotulos_linhas, key_deslocamentos_linhas
):
    """
    Renderiza um mapa de correspondência completo (estático + interativo)
    para o par (var_linha, var_coluna).

    Os rótulos/deslocamentos das LINHAS vêm de fora (key_rotulos_linhas,
    key_deslocamentos_linhas) — são compartilhados entre todos os mapas,
    pois a variável principal é sempre a mesma nas linhas. Já os rótulos
    das COLUNAS são específicos deste par, isolados por 'key_prefix'.

    Não altera nem é chamada por `analise_correspondencia`.
    """
    try:
        coordenadas_linhas, coordenadas_colunas, eigenvalues = _crosstab_e_ca(
            dados, var_linha, var_coluna, 2
        )
    except Exception as e:
        st.error(
            f"Não foi possível calcular a Análise de Correspondência para "
            f"'{var_linha}' x '{var_coluna}': {e}"
        )
        return

    if coordenadas_linhas.shape[1] < 2 or coordenadas_colunas.shape[1] < 2:
        st.error(
            f"A análise de correspondência para '{var_linha}' x '{var_coluna}' "
            "não conseguiu gerar duas dimensões. Verifique os dados selecionados."
        )
        return

    st.subheader("Inércia explicada (variância)")
    explained_inertia = eigenvalues / eigenvalues.sum()
    for i, valor in enumerate(explained_inertia):
        st.write(f"Dim {i+1}: {valor * 100:.2f}%")

    # garante que toda categoria de linha desta tabela específica tenha
    # entrada no dicionário compartilhado (ex.: categoria nova, sem edição ainda)
    for i in coordenadas_linhas.index:
        st.session_state[key_rotulos_linhas].setdefault(i, str(i))
        st.session_state[key_deslocamentos_linhas].setdefault(i, (0.0, 0.0))

    key_rotulos_colunas = f"rotulos_colunas__{key_prefix}"
    key_deslocamentos_colunas = f"deslocamentos_colunas__{key_prefix}"

    if key_rotulos_colunas not in st.session_state:
        st.session_state[key_rotulos_colunas] = {
            i: str(i) for i in coordenadas_colunas.index
        }
    else:
        for i in coordenadas_colunas.index:
            st.session_state[key_rotulos_colunas].setdefault(i, str(i))

    if key_deslocamentos_colunas not in st.session_state:
        st.session_state[key_deslocamentos_colunas] = {
            i: (0.0, 0.0) for i in coordenadas_colunas.index
        }
    else:
        for i in coordenadas_colunas.index:
            st.session_state[key_deslocamentos_colunas].setdefault(i, (0.0, 0.0))

    legenda_linhas = st.text_input(
        "Legenda para Linhas:", var_linha, key=f"legenda_linhas__{key_prefix}"
    )
    legenda_colunas = st.text_input(
        "Legenda para Colunas:", var_coluna, key=f"legenda_colunas__{key_prefix}"
    )
    mostrar_legenda = st.checkbox(
        "Mostrar legenda (Linhas/Colunas)", value=False, key=f"mostrar_legenda__{key_prefix}"
    )

    with st.expander(f"Editar categorias de '{var_coluna}' (colunas deste mapa)"):
        with st.form(key=f"form_colunas__{key_prefix}"):
            rotulo_tmp = {}
            dx_tmp = {}
            dy_tmp = {}
            remover_tmp = {}

            for i in coordenadas_colunas.index:
                col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
                rotulo_atual = st.session_state[key_rotulos_colunas][i]
                dx_atual, dy_atual = st.session_state[key_deslocamentos_colunas][i]

                with col1:
                    rotulo_tmp[i] = st.text_input(
                        f"Rótulo — '{i}'",
                        value=rotulo_atual if rotulo_atual is not None else str(i),
                        key=f"rotulo_coluna__{key_prefix}__{i}"
                    )
                with col2:
                    dx_tmp[i] = st.number_input(
                        "Δx", min_value=-1.0, max_value=1.0,
                        value=float(dx_atual), step=0.01, format="%.2f",
                        key=f"desloc_x_coluna__{key_prefix}__{i}"
                    )
                with col3:
                    dy_tmp[i] = st.number_input(
                        "Δy", min_value=-1.0, max_value=1.0,
                        value=float(dy_atual), step=0.01, format="%.2f",
                        key=f"desloc_y_coluna__{key_prefix}__{i}"
                    )
                with col4:
                    remover_tmp[i] = st.checkbox(
                        "Remover",
                        value=(rotulo_atual is None),
                        key=f"remover_coluna__{key_prefix}__{i}"
                    )

            aplicar = st.form_submit_button("Aplicar neste mapa")

            if aplicar:
                for i in coordenadas_colunas.index:
                    st.session_state[key_rotulos_colunas][i] = (
                        None if remover_tmp[i] else rotulo_tmp[i]
                    )
                    st.session_state[key_deslocamentos_colunas][i] = (dx_tmp[i], dy_tmp[i])
                st.success("Alterações aplicadas neste mapa.")

    max_deslocamento = 0.05
    pontos_linhas = []
    pontos_colunas = []

    for i, row in coordenadas_linhas.iterrows():
        if st.session_state[key_rotulos_linhas][i] is None:
            continue
        x = row.iloc[0]
        y = row.iloc[1]
        dx, dy = st.session_state[key_deslocamentos_linhas][i]
        dx = max(-max_deslocamento, min(max_deslocamento, dx * (x / abs(x) if x != 0 else 1)))
        dy = max(-max_deslocamento, min(max_deslocamento, dy * (y / abs(y) if y != 0 else 1)))
        pontos_linhas.append({
            "x": x,
            "y": y,
            "dx": dx,
            "dy": dy,
            "label": st.session_state[key_rotulos_linhas][i]
        })

    for i, row in coordenadas_colunas.iterrows():
        if st.session_state[key_rotulos_colunas][i] is None:
            continue
        x = row.iloc[0]
        y = row.iloc[1]
        dx, dy = st.session_state[key_deslocamentos_colunas][i]
        dx = max(-max_deslocamento, min(max_deslocamento, dx * (x / abs(x) if x != 0 else 1)))
        dy = max(-max_deslocamento, min(max_deslocamento, dy * (y / abs(y) if y != 0 else 1)))
        pontos_colunas.append({
            "x": x,
            "y": y,
            "dx": dx,
            "dy": dy,
            "label": st.session_state[key_rotulos_colunas][i]
        })

    st.subheader("Mapa estático")

    fig, ax = plt.subplots(figsize=(8, 6))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    texts = []
    linha_legenda_plotada = False
    coluna_legenda_plotada = False

    for p in pontos_linhas:
        label_legenda = legenda_linhas if (mostrar_legenda and not linha_legenda_plotada) else ""
        ax.scatter(p["x"], p["y"], color="blue", marker="o", s=20, label=label_legenda)
        if mostrar_legenda and not linha_legenda_plotada:
            linha_legenda_plotada = True

        txt = ax.text(
            p["x"] + p["dx"],
            p["y"] + p["dy"],
            p["label"],
            color="blue",
            fontsize=8,
            ha="center",
            va="bottom",
            fontweight="bold",
        )
        texts.append(txt)

    for p in pontos_colunas:
        label_legenda = legenda_colunas if (mostrar_legenda and not coluna_legenda_plotada) else ""
        ax.scatter(p["x"], p["y"], color="red", marker="^", s=30, label=label_legenda)
        if mostrar_legenda and not coluna_legenda_plotada:
            coluna_legenda_plotada = True

        txt = ax.text(
            p["x"] + p["dx"],
            p["y"] + p["dy"],
            p["label"],
            color="red",
            fontsize=8,
            ha="center",
            va="bottom",
            fontweight="bold",
        )
        texts.append(txt)

    adjust_text(
        texts,
        arrowprops=None,
        force_text=(0.5, 1),
        force_points=(0.5, 1),
        expand_text=(1.2, 1.5),
        expand_points=(1.2, 1.5),
        lim=100,
    )

    ax.axhline(0, color="black", linestyle="--", linewidth=0.8)
    ax.axvline(0, color="black", linestyle="--", linewidth=0.8)
    ax.grid(color="lightgray", linestyle="--", linewidth=0.5, alpha=0.7)
    ax.set_aspect("auto")

    if mostrar_legenda and (linha_legenda_plotada or coluna_legenda_plotada):
        ax.legend(loc="upper right", frameon=False, fontsize=12)

    ax.set_title(f"{var_linha} x {var_coluna}", fontsize=14)
    st.pyplot(fig)

    with st.expander("Mapa interativo (clique para exibir)", expanded=False):
        fig_int = go.Figure()

        fig_int.add_trace(go.Scatter(
            x=[p["x"] for p in pontos_linhas],
            y=[p["y"] for p in pontos_linhas],
            mode="markers",
            name=legenda_linhas,
            marker=dict(color="blue", symbol="circle", size=8),
            showlegend=mostrar_legenda and len(pontos_linhas) > 0
        ))

        fig_int.add_trace(go.Scatter(
            x=[p["x"] for p in pontos_colunas],
            y=[p["y"] for p in pontos_colunas],
            mode="markers",
            name=legenda_colunas,
            marker=dict(color="red", symbol="triangle-up", size=10),
            showlegend=mostrar_legenda and len(pontos_colunas) > 0
        ))

        annotations = []
        for p in pontos_linhas:
            annotations.append(dict(
                x=p["x"] + p["dx"],
                y=p["y"] + p["dy"],
                xref="x",
                yref="y",
                text=p["label"],
                showarrow=False,
                font=dict(color="blue", size=10),
                xanchor="center",
                yanchor="bottom",
            ))
        for p in pontos_colunas:
            annotations.append(dict(
                x=p["x"] + p["dx"],
                y=p["y"] + p["dy"],
                xref="x",
                yref="y",
                text=p["label"],
                showarrow=False,
                font=dict(color="red", size=10),
                xanchor="center",
                yanchor="bottom",
            ))

        xs = [p["x"] for p in pontos_linhas] + [p["x"] for p in pontos_colunas]
        ys = [p["y"] for p in pontos_linhas] + [p["y"] for p in pontos_colunas]

        if xs and ys:
            x_min, x_max = min(xs), max(xs)
            y_min, y_max = min(ys), max(ys)
            dx_range = (x_max - x_min) * 0.1 or 0.1
            dy_range = (y_max - y_min) * 0.1 or 0.1
            x_range = [x_min - dx_range, x_max + dx_range]
            y_range = [y_min - dy_range, y_max + dy_range]
        else:
            x_range = [-1, 1]
            y_range = [-1, 1]

        fig_int.update_layout(
            template=None,
            title=dict(
                text=f"{var_linha} x {var_coluna}",
                font=dict(color="black", size=14),
                x=0.5
            ),
            font=dict(color="black"),
            paper_bgcolor="white",
            plot_bgcolor="white",
            annotations=annotations,
            showlegend=mostrar_legenda,
            legend=dict(
                x=0.99, y=0.99,
                xanchor="right", yanchor="top",
                bgcolor="rgba(0,0,0,0)",
                font=dict(color="black", size=12),
            ),
            margin=dict(l=40, r=20, t=60, b=40),
            xaxis=dict(range=x_range),
            yaxis=dict(range=y_range),
        )

        fig_int.update_xaxes(
            title=dict(text="", font=dict(color="black")),
            showgrid=True,
            gridcolor="lightgray",
            gridwidth=0.5,
            griddash="dash",
            zeroline=False,
            tickfont=dict(color="black"),
            linecolor="black",
            mirror=True,
        )
        fig_int.update_yaxes(
            title=dict(text="", font=dict(color="black")),
            showgrid=True,
            gridcolor="lightgray",
            gridwidth=0.5,
            griddash="dash",
            zeroline=False,
            tickfont=dict(color="black"),
            linecolor="black",
            mirror=True,
        )

        fig_int.add_shape(
            type="line",
            x0=x_range[0], x1=x_range[1],
            y0=0, y1=0,
            line=dict(color="black", width=1, dash="dash"),
            xref="x", yref="y",
            layer="above"
        )
        fig_int.add_shape(
            type="line",
            x0=0, x1=0,
            y0=y_range[0], y1=y_range[1],
            line=dict(color="black", width=1, dash="dash"),
            xref="x", yref="y",
            layer="above"
        )

        config = {
            "editable": True,
            "edits": {"annotationPosition": True},
        }

        st.plotly_chart(
            fig_int,
            use_container_width=True,
            config=config,
            key=f"plotly_chart__{key_prefix}"
        )

    return fig


# =========================================================
# Exportação dos mapas estáticos para PowerPoint
# =========================================================
def _gerar_pptx_mapas(figuras):
    """
    Recebe uma lista de tuplas (titulo, fig_matplotlib), na mesma ordem em
    que os mapas foram gerados na tela, e monta uma apresentação
    PowerPoint com um slide por mapa (título + imagem do mapa estático).
    Retorna um BytesIO pronto para download.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout_em_branco = prs.slide_layouts[6]

    for titulo, fig in figuras:
        slide = prs.slides.add_slide(layout_em_branco)

        caixa_titulo = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)
        )
        tf = caixa_titulo.text_frame
        tf.text = titulo
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(1.8), Inches(1.3), height=Inches(5.7))

    saida = io.BytesIO()
    prs.save(saida)
    saida.seek(0)
    return saida


# =========================================================
# Mapas de Correspondência - Múltiplas Variáveis
# =========================================================
def analise_correspondencia_multipla(dados):
    """
    Permite escolher uma variável principal (ex.: P7) e marcar de uma vez
    várias outras variáveis (ex.: P1, P2, P3...) para cruzar com ela.
    Gera um mapa de correspondência completo e editável para cada par
    (variável principal x variável secundária).

    Se o arquivo carregado tiver colunas que batem com
    VARIAVEIS_INTERESSE_PADRAO, o app já sugere/pré-seleciona essas
    variáveis, para agilizar a geração dos mapas.

    Não modifica e não é chamada por `analise_correspondencia`.
    """
    st.header("Mapas de Correspondência — Múltiplas Variáveis")

    colunas = dados.columns.tolist()

    disponiveis = [v for v in VARIAVEIS_INTERESSE_PADRAO if v in colunas]
    faltantes = [v for v in VARIAVEIS_INTERESSE_PADRAO if v not in colunas]

    if disponiveis:
        st.success(
            f"Encontradas {len(disponiveis)} de {len(VARIAVEIS_INTERESSE_PADRAO)} "
            f"variáveis de interesse neste arquivo: {', '.join(disponiveis)}."
        )
    if faltantes:
        st.caption(
            f"Não encontradas neste arquivo: {', '.join(faltantes)}."
        )

    # sugere 'P7' (ou a primeira variável de interesse disponível) como principal
    if "P7" in disponiveis:
        principal_sugerida = "P7"
    elif disponiveis:
        principal_sugerida = disponiveis[0]
    else:
        principal_sugerida = colunas[0]

    variavel_principal = st.selectbox(
        "Selecione a variável principal:",
        colunas,
        index=colunas.index(principal_sugerida),
        key="ca_multi_var_principal"
    )

    colunas_restantes = [c for c in colunas if c != variavel_principal]

    # sugere as demais variáveis de interesse (exceto a principal) como
    # padrão do multiselect; a key muda junto com a principal para que a
    # sugestão se atualize se o usuário trocar a variável principal
    sugestao_secundarias = [v for v in disponiveis if v != variavel_principal]

    variaveis_secundarias = st.multiselect(
        f"Selecione as variáveis para cruzar com '{variavel_principal}':",
        colunas_restantes,
        default=sugestao_secundarias,
        key=f"ca_multi_var_secundarias__{variavel_principal}"
    )

    if not variaveis_secundarias:
        st.warning(
            "Selecione pelo menos uma variável para cruzar com a variável principal."
        )
        return

    st.caption(
        f"{len(variaveis_secundarias)} mapa(s) será(ão) gerado(s) para "
        f"'{variavel_principal}' x cada variável selecionada."
    )

    # configuração das categorias da variável principal (linhas),
    # compartilhada por todos os mapas gerados abaixo
    key_rotulos_linhas, key_deslocamentos_linhas = _configurar_categorias_principal(
        dados, variavel_principal
    )

    figuras_geradas = []

    for var_secundaria in variaveis_secundarias:
        key_prefix = f"{variavel_principal}__{var_secundaria}"
        titulo_mapa = f"{variavel_principal} x {var_secundaria}"
        with st.expander(titulo_mapa, expanded=True):
            fig = _renderizar_mapa_correspondencia(
                dados,
                variavel_principal,
                var_secundaria,
                key_prefix,
                key_rotulos_linhas,
                key_deslocamentos_linhas,
            )
        if fig is not None:
            figuras_geradas.append((var_secundaria, fig))

    if figuras_geradas:
        st.markdown("---")
        st.subheader("Exportar mapas")

        if not PPTX_AVAILABLE:
            st.warning(
                "A biblioteca 'python-pptx' não está instalada neste ambiente "
                f"({PPTX_IMPORT_ERROR}). Adicione 'python-pptx' ao "
                "requirements.txt para habilitar a exportação em PowerPoint."
            )
        else:
            st.caption(
                f"{len(figuras_geradas)} mapa(s) estático(s) serão exportados, "
                "um por slide, na mesma ordem em que foram gerados acima."
            )

            titulo_principal_pptx = st.text_input(
                "Qual título você quer colocar na variável principal "
                f"(variável selecionada no início do app, '{variavel_principal}')?",
                value=_rotulo_amigavel(variavel_principal),
                key=f"titulo_pptx_principal__{variavel_principal}"
            )

            if st.button("Gerar PowerPoint com todos os mapas", key="gerar_pptx_multi"):
                figuras_com_titulo = [
                    (f"{titulo_principal_pptx} x {_rotulo_amigavel(var_secundaria)}", fig)
                    for var_secundaria, fig in figuras_geradas
                ]
                with st.spinner("Montando a apresentação..."):
                    pptx_buffer = _gerar_pptx_mapas(figuras_com_titulo)
                st.session_state["pptx_mapas_bytes"] = pptx_buffer.getvalue()
                st.success("PowerPoint gerado! Use o botão abaixo para baixar.")

            if "pptx_mapas_bytes" in st.session_state:
                st.download_button(
                    label="Baixar PowerPoint",
                    data=st.session_state["pptx_mapas_bytes"],
                    file_name=f"mapas_correspondencia_{variavel_principal}.pptx",
                    mime=(
                        "application/vnd.openxmlformats-officedocument"
                        ".presentationml.presentation"
                    ),
                    key="download_pptx_multi"
                )


# =========================
# Função principal
# =========================
def main():
    uploaded_file = st.file_uploader("Carregar Arquivo", type=["csv", "xlsx", "sav"])

    if uploaded_file is not None:
        dados, dados_exibicao = carregar_dados(uploaded_file)

        if dados is not None:
            analise = st.sidebar.selectbox(
                "Escolha a análise",
                (
                    "Tratamento de Dados",
                    "Mapas de Correspondência",
                    "Mapas de Correspondência (Múltiplas Variáveis)",
                ),
            )

            if analise == "Mapas de Correspondência":
                analise_correspondencia(dados)

            elif analise == "Mapas de Correspondência (Múltiplas Variáveis)":
                analise_correspondencia_multipla(dados)

    else:
        st.info("Por favor, carregue um arquivo para começar.")


if __name__ == "__main__":
    main()
